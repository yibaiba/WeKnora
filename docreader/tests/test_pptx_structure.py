from io import BytesIO
import unittest

from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches

from function_test_loader import load_function_tests

from docreader.models.document import Document, StructureBlock
from docreader.parser.base_parser import BaseParser
from docreader.parser.chain_parser import PipelineParser
from docreader.parser.pptx_structure import attach_pptx_structure
from docreader.parser.structure_relocation import relocate_structure_blocks
from docreader.structure import structure_blocks_from_document


def _structured_pptx() -> bytes:
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[5])
    first.shapes.title.text = "Quarterly Review"
    text_box = first.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
    paragraph = text_box.text_frame.paragraphs[0]
    paragraph.text = "Stable throughput"
    bullet = OxmlElement("a:buChar")
    bullet.set("char", "•")
    paragraph._p.get_or_add_pPr().append(bullet)
    table = first.shapes.add_table(
        2, 2, Inches(1), Inches(4), Inches(6), Inches(1.2)
    ).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "State"
    table.cell(1, 0).text = "Recall"
    table.cell(1, 1).text = "Good"

    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "Actions"
    action = second.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
    action.text_frame.paragraphs[0].text = "Monitor the next release"
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _final_markdown() -> str:
    return (
        "<!-- Slide number: 1 -->\n"
        "# Quarterly Review\n"
        "Stable throughput\n"
        "| Metric | State |\n"
        "| --- | --- |\n"
        "| Recall | Good |\n\n"
        "<!-- Slide number: 2 -->\n"
        "# Actions\n"
        "Monitor the next release\n"
    )


def test_pptx_native_structure_maps_slides_text_boxes_tables_and_transport():
    markdown = _final_markdown()
    document = attach_pptx_structure(
        _structured_pptx(), Document(content=markdown)
    )

    kinds = [block.kind for block in document.structure_blocks]
    assert kinds == [
        "page_region",
        "heading",
        "list_item",
        "table_header",
        "table_row",
        "page_region",
        "heading",
        "paragraph",
    ]
    assert all(markdown[block.start : block.end] for block in document.structure_blocks)
    assert len({block.id for block in document.structure_blocks}) == len(
        document.structure_blocks
    )
    first_heading = document.structure_blocks[1]
    assert all(
        block.parent_id == first_heading.id
        for block in document.structure_blocks[2:5]
    )
    assert document.structure_blocks[3].table_id == document.structure_blocks[4].table_id
    assert document.structure_blocks[-1].parent_id == document.structure_blocks[-2].id
    transported = structure_blocks_from_document(document)
    assert [block.kind for block in transported] == kinds


def test_invalid_pptx_metadata_preserves_markdown_and_exposes_reason():
    markdown = "Preserved presentation body"

    document = attach_pptx_structure(b"not-a-pptx", Document(content=markdown))

    assert document.content == markdown
    assert document.structure_blocks == []
    assert document.metadata == {"semantic_structure_status": "pptx_metadata_invalid"}
    assert markdown not in document.metadata["semantic_structure_status"]


def test_markitdown_pptx_production_path_emits_final_structure():
    try:
        from docreader.parser.markitdown_parser import MarkitdownParser
    except ImportError as error:
        raise unittest.SkipTest("MarkItDown runtime dependencies unavailable") from error

    document = MarkitdownParser(
        file_name="quarterly-review.pptx", file_type="pptx"
    ).parse_into_text(_structured_pptx())

    assert document.content
    assert {"page_region", "heading", "table_header", "table_row"}.issubset(
        {block.kind for block in document.structure_blocks}
    )
    assert all(
        document.content[block.start : block.end]
        for block in document.structure_blocks
    )


class _StructuredStage(BaseParser):
    def parse_into_text(self, content: bytes) -> Document:
        markdown = "# Title\n\n![Chart](images/chart.png)\n\nBody\n"
        return Document(
            content=markdown,
            images={"images/chart.png": "payload"},
            metadata={"source": "first"},
            structure_blocks=[
                _block(markdown, "heading", "# Title\n", "heading-1"),
                _block(
                    markdown,
                    "image_caption",
                    "![Chart](images/chart.png)\n",
                    "image-1",
                    parent_id="heading-1",
                    context_kinds=["image"],
                ),
                _block(
                    markdown,
                    "paragraph",
                    "Body\n",
                    "body-1",
                    parent_id="heading-1",
                ),
            ],
        )


class _RewriteStage(BaseParser):
    def parse_into_text(self, content: bytes) -> Document:
        return Document(
            content=(
                "\n# Title\n\n![Chart](stored/chart-long-name.png)\n\nBody  \n"
            ),
            images={"stored/chart-long-name.png": "stored"},
            metadata={"stage": "rewrite"},
        )


def _block(
    markdown: str,
    kind: str,
    needle: str,
    block_id: str,
    *,
    parent_id: str = "",
    context_kinds: list[str] | None = None,
) -> StructureBlock:
    start = markdown.index(needle)
    return StructureBlock(
        id=block_id,
        kind=kind,
        start=start,
        end=start + len(needle),
        parent_id=parent_id,
        atomic=kind != "paragraph",
        confidence="high",
        context_kinds=context_kinds or [],
        section_depth=1 if kind == "heading" else 0,
    )


def test_pipeline_relocates_upstream_hints_and_merges_images_metadata():
    parser_type = PipelineParser.create(_StructuredStage, _RewriteStage)

    with unittest.TestCase().assertLogs(
        "docreader.parser.chain_parser", level="WARNING"
    ) as captured:
        document = parser_type(file_type="pptx").parse_into_text(b"presentation")

    assert [block.id for block in document.structure_blocks] == [
        "heading-1",
        "body-1",
    ]
    assert all(
        document.content[block.start : block.end].strip()
        for block in document.structure_blocks
    )
    assert document.structure_blocks[1].parent_id == "heading-1"
    assert document.images == {
        "images/chart.png": "payload",
        "stored/chart-long-name.png": "stored",
    }
    assert document.metadata == {"source": "first", "stage": "rewrite"}
    assert "pipeline_hint_source_unmatched" in "\n".join(captured.output)


def test_relocation_rejects_ambiguous_normalized_text():
    source = "Repeated value\n"
    block = _block(source, "paragraph", source, "paragraph-1")

    result = relocate_structure_blocks(
        source,
        "Repeated   value\n\nRepeated\tvalue\n",
        [block],
    )

    assert result.blocks == ()
    assert result.rejected_count == 1
    assert result.reason_codes == ("pipeline_hint_normalized_ambiguous",)


def load_tests(loader, tests, pattern):
    del loader, pattern
    return load_function_tests(globals(), tests)
